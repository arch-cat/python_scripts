# -*- coding: utf-8 -*-
# Imports
import io
import os
import sys
import warnings
import logging
import paramiko
import colorama
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

import matplotlib.dates as plt_dates
import matplotlib.pyplot as plt
import pandas as pd
import win32com.client as win32  # Библитека для работы старой функции send_msg_old()

import settings
import logger

warnings.filterwarnings("ignore")
colorama.init()


def watch_logs_dau_mau_wau():
    command = "./parse_dau_mau_wau.sh"
    logging.info("SSH connection...")
    with paramiko.SSHClient() as client:
        start_time = datetime.now()
        client.load_system_host_keys()
        client.connect(
            hostname=settings.SSH_HOST,
            username=settings.SSH_USER,
            password=settings.SSH_PASSWORD,
        )
        logging.info(f"Successfully logged in as {settings.SSH_USER}")
        stdout_list = []
        (stdin, stdout, stderr) = client.exec_command(command, get_pty=True)
        for line in iter(stdout.readline, ""):
            print(line, end="")
            stdout_list.append(line)
    logging.info("SSH connect close...")
    stdout = "".join(stdout_list)
    parse_data = stdout[stdout.find("Вывод данных в терминал:") + 24 :]
    return start_time, parse_data


def create_plt_bar_dau(_parse_data):
    logging.info(f"Генерируем график для дней")
    dau = _parse_data[_parse_data.find("dau") + 3 : _parse_data.find("wau")]
    convert_io = io.StringIO(dau)
    df = pd.read_csv(
        convert_io, delimiter=" ", index_col=["date"], parse_dates=["date"]
    )
    plt.style.use("seaborn")

    fig, ax = plt.subplots(figsize=(16.3, 3))
    ax.bar(df.index, df["count"], width=0.36)  # width=0.32, 0.34
    ax.xaxis.set_major_locator(plt_dates.MonthLocator())
    ax.xaxis.set_major_formatter(plt_dates.DateFormatter("%Y-%m"))
    fig.autofmt_xdate()
    ax.set_xlim([df.index.min(), df.index.max()])
    ax.set_title("Dau")

    if os.path.isfile("./tmp/dau.png"):
        os.remove("./tmp/dau.png")
    plt.savefig("./tmp/dau.png", bbox_inches="tight")
    logging.info(f"График успешно создан")


def create_plt_bar_wau(_parse_data):
    logging.info(f"Генерируем график для недель")
    wau = _parse_data[_parse_data.find("wau") + 3 : _parse_data.find("mau")]
    convert_io = io.StringIO(wau)
    df = pd.read_csv(convert_io, delimiter=" ")

    plt.rcParams["figure.figsize"] = [10.2, 3.5]

    df.plot.bar(x="date", y="count", fontsize=8, legend=None)
    plt.ylim(0, 1400)
    plt.xlabel("")
    plt.title("Wau")
    plt.style.use("seaborn")

    for i, val in enumerate(df["count"].values):
        plt.text(
            i,
            val,
            ("  " + str(val)),
            rotation="vertical",
            horizontalalignment="center",
            verticalalignment="bottom",
            fontsize=8,
        )

    if os.path.isfile("./tmp/wau.png"):
        os.remove("./tmp/wau.png")
    plt.savefig("./tmp/wau.png", bbox_inches="tight")
    logging.info(f"График успешно создан")


def create_plt_bar_mau(_parse_data):
    logging.info(f"Генерируем график для месяцев")
    mau = _parse_data[_parse_data.find("mau") + 4 :]
    convert_io = io.StringIO(mau)
    df = pd.read_csv(convert_io, delimiter=" ")

    plt.rcParams["figure.figsize"] = [5.5, 3.6]

    df.plot.bar(x="date", y="count", fontsize=9, legend=None)
    plt.ylim(0, 2500)
    plt.xlabel("")
    plt.title("Mau")
    plt.style.use("seaborn")

    for i, val in enumerate(df["count"].values):
        plt.text(
            i,
            val,
            ("  " + str(val)),
            rotation="vertical",
            horizontalalignment="center",
            verticalalignment="bottom",
            fontsize=8,
        )

    if os.path.isfile("./tmp/mau.png"):
        os.remove("./tmp/mau.png")
    plt.savefig("./tmp/mau.png", bbox_inches="tight")
    logging.info(f"График успешно создан")


def create_pptx_presentation():
    get_current_date = datetime.now().strftime("%d-%m-%Y")
    logging.info(f"Создаем файл презентации с графиками")
    img_path_dau = r"./tmp/dau.png"
    img_path_mau = r"./tmp/mau.png"
    img_path_wau = r"./tmp/wau.png"

    presentation_gp = open(
        r"./tmp/КАП - Шаблон.pptx",
        "rb",
    )
    ppt = Presentation(presentation_gp)

    slide = ppt.slides[0]
    slide.shapes.add_picture(img_path_dau, Inches(0), Inches(0.7))
    slide.shapes.add_picture(img_path_wau, Inches(0), Inches(3.7))
    slide.shapes.add_picture(img_path_mau, Inches(8.48), Inches(3.74))

    pptx = f"./pptx/Шаблон - {get_current_date}.pptx"
    ppt.save(pptx)
    presentation_gp.close()
    logging.info(f"Презентация успешно создана")
    return pptx


def send_msg_old(_save_pptx, _end_time):
    logging.info(f"Направляем письмо с презентацией")
    get_current_date = datetime.now().strftime("%d-%m-%Y")
    outlook = win32.Dispatch("outlook.application")
    mail = outlook.CreateItem(0)
    # Email destination "Email#1;Email#2"
    mail.To = ""
    mail.Subject = f"Статус {get_current_date}"
    mail.HTMLBody = (
        f"<p>Сформировано автоматически</p>"
        f"Время выполнения {_end_time}</p>"
        f"<p>-----------------------------------------------------<br>"
        f"(C) 2023 Grif Ivan, Moscow, Russia<br>"
        f"Released under GNU Public License (GPL)<br>"
        f"<a href=https://bitbucket.ru/projects/repos/"
        f"browse/support/grif_iv?at=support>Bitbucket</a></p>"
    )
    mail.Attachments.Add(os.path.join(os.getcwd(), _save_pptx))
    mail.Send()
    logging.info(f"Письмо направлено успешно")


def send_msg(_save_pptx, _end_time):
    logging.info(f"Направляем письмо с презентацией")
    get_current_date = datetime.now().strftime("%d-%m-%Y")
    user = ""
    server = smtplib.SMTP("", 587)
    server.starttls()
    server.login(user, "")
    msg = MIMEMultipart()
    msg["From"] = f""
    msg["To"] = f""
    msg["Subject"] = f"Статус {get_current_date}"
    msg.attach(
        MIMEText(
            f"""
        <html>
            <body>
                <p>Сформировано автоматически</p>
                <p>Время выполнения {_end_time}</p>
                <p>-----------------------------------------------------<br>
                (C) 2023 Grif Ivan, Moscow, Russia<br>
                Released under GNU Public License (GPL)<br>
                <a href=https://bitbucket.ru/projects/repos/browse/support\
                /grif_iv?at=support>Bitbucket</a></p> 
            </body>
        </html>
        """,
            "html",
            "utf-8",
        )
    )
    pptx = MIMEApplication(open(f"{_save_pptx}", "rb").read())
    pptx.add_header(
        "Content-Disposition",
        "attachment",
        filename=f"{_save_pptx.partition('-')[2].strip()}",
    )
    msg.attach(pptx)
    server.send_message(msg)
    server.close()
    logging.info(f"Письмо направлено успешно")


def main():
    get_current_date = datetime.now().strftime("%d-%m-%Y")
    script_name = os.path.splitext(os.path.basename(sys.argv[0]))[0]
    days = (
        "Monday",
        "Tuesday",
        "Wednesday",
        "Thursday",
        "Friday",
        "Saturday",
        "Sunday",
    )
    now_dt = datetime.now()
    current_weekday = days[now_dt.weekday()]
    if not logger.set_up_logging(
        console_log_output="stdout",
        console_log_level="info",
        console_log_color=True,
        logfile_file=f"./logs/{script_name}_{get_current_date}.log",
        logfile_log_level="debug",
        logfile_log_color=False,
        log_line_template=f"%(color_on)s[{current_weekday} %(asctime)s] [%(threadName)s] [%(levelname)s] :::: %("
        f"message)s%(color_off)s",
    ):
        print("Failed to set up logging, aborting.")
        return 1

    logging.info(f"Begin {script_name} script")
    _start_time, _parse_data = watch_logs_dau_mau_wau()
    create_plt_bar_dau(_parse_data)
    create_plt_bar_wau(_parse_data)
    create_plt_bar_mau(_parse_data)
    _save_pptx = create_pptx_presentation()
    _end_time = datetime.now() - _start_time
    _end_time = str(_end_time).split(".")[0]
    send_msg(_save_pptx, _end_time)
    logging.info(f"End {script_name} script")
